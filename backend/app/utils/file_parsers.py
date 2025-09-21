from __future__ import annotations

import csv
import io
import json
from pathlib import Path
from typing import Callable, Dict

from pdfminer.high_level import extract_text as pdf_extract_text
from docx import Document
from pptx import Presentation
import openpyxl


EXTENSION_MAP: Dict[str, Callable[[Path, bytes], str]] = {}


def register(extension: str) -> Callable[[Callable[[Path, bytes], str]], Callable[[Path, bytes], str]]:
    def decorator(func: Callable[[Path, bytes], str]) -> Callable[[Path, bytes], str]:
        EXTENSION_MAP[extension.lower()] = func
        return func

    return decorator


@register(".txt")
def _parse_text(path: Path, raw: bytes) -> str:
    return raw.decode("utf-8", errors="ignore")


@register(".md")
def _parse_markdown(path: Path, raw: bytes) -> str:
    return raw.decode("utf-8", errors="ignore")


@register(".pdf")
def _parse_pdf(path: Path, raw: bytes) -> str:
    try:
        return pdf_extract_text(str(path))
    except Exception:
        return ""


@register(".docx")
def _parse_docx(path: Path, raw: bytes) -> str:
    try:
        doc = Document(str(path))
    except Exception:
        return ""
    return "\n".join(paragraph.text for paragraph in doc.paragraphs)


@register(".pptx")
def _parse_pptx(path: Path, raw: bytes) -> str:
    try:
        presentation = Presentation(str(path))
    except Exception:
        return ""
    slides = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slides.append(shape.text)
    return "\n".join(slides)


@register(".csv")
def _parse_csv(path: Path, raw: bytes) -> str:
    stream = io.StringIO(raw.decode("utf-8", errors="ignore"))
    reader = csv.reader(stream)
    rows = [", ".join(row) for row in reader]
    return "\n".join(rows)


@register(".json")
def _parse_json(path: Path, raw: bytes) -> str:
    try:
        data = json.loads(raw.decode("utf-8", errors="ignore"))
        return json.dumps(data, indent=2)
    except Exception:
        return ""


@register(".xlsx")
def _parse_xlsx(path: Path, raw: bytes) -> str:
    wb = openpyxl.load_workbook(str(path), data_only=True)
    contents = []
    for sheet in wb.worksheets:
        contents.append(f"# Sheet: {sheet.title}")
        for row in sheet.iter_rows(values_only=True):
            row_values = [str(cell) for cell in row if cell is not None]
            if row_values:
                contents.append(", ".join(row_values))
    return "\n".join(contents)


def extract_text_from_file(path: Path, raw: bytes) -> str:
    parser = EXTENSION_MAP.get(path.suffix.lower())
    if parser:
        return parser(path, raw)
    return raw.decode("utf-8", errors="ignore")
